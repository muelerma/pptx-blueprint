import pathlib
import pptx
from typing import Union, Iterable, Tuple
from pptx.shapes.base import BaseShape

_Pathlike = Union[str, pathlib.Path]


class Template:
    """Helper class for modifying pptx templates.
    """

    def __init__(self, filename: _Pathlike) -> None:
        """Initializes a Template-Modifier.

        Args:
            filename (path-like): file name or path
        """
        self._template_path = filename
        self._presentation = pptx.Presentation(filename)
        pass

    def replace_text(self, label: str, new_text: str) -> None:
        """Replaces text placeholders on one or many slides.

        Args:
            label (str): label of the placeholder (without curly braces)
            text (str): new content
            scope: None, slide number, Slide object or iterable of Slide objects
        """
        slide_number, tag_name = self._parse_label(label)
        shapes = self._find_shapes(slide_number, tag_name)
        for shape in shapes:
            shape.text = new_text

    def replace_picture(self, label: str, filename: _Pathlike) -> None:
        """Replaces rectangle placeholders on one or many slides.

        Args:
            label (str): label of the placeholder (without curly braces)
            filename (path-like): path to an image file
        """
        pass

    def replace_table(self, label: str, data, header=False, rownames=False):
        """Replaces rectangle placeholders on one or many slides.

        Args:
            label (str): label of the placeholder (without curly braces)
            data (pandas.DataFrame): table to be inserted into the presentation
        """
        assert isinstance(data, pandas.dataframe)
        
        shapes_to_replace = self._find_shapes(label)

        rows, cols = data.shape
        if header: rows += 1
        if rownames: cols += 1

        for old_shape in shapes_to_replace:
            slide_shapes = old_shape._parent
            table = slide_shapes.add_table(
                rows,
                cols,
                old_shape.left, 
                old_shape.top, 
                old_shape.width,
                old_shape.height
            ).table
            
            # set column widths
            col_width = Length(old_shape.width/cols)
            for i in range(len(table.columns)):
                table.columns[i].width = col_width

            ## fill the table
            for c in range(data.shape[1]):
                for r in range(data.shape[0]):
                    if header and r == 0:
                        c_temp = c
                        ## when rownames, skip first column
                        if rownames: c_temp += 1
                        table.cell(r, c_temp).text = data.columns[c]
                    if rownames and c == 0:
                        r_temp = r
                        ## when header, skip first row
                        if header: r_temp += 1
                        table.cell(r_temp, c).text = str(data.index[r])
                    ## fill table body
                    r_shape = r
                    if header: r_shape += 1
                    c_shape = c
                    if rownames: c_shape += 1
                    table.cell(r_shape, c_shape).text = str(data.iloc[r, c])


    def _parse_label(self, label: str) -> Tuple[Union[int, str], str]:
        slide_number, tag_name = label.split(':')
        return int(slide_number) if slide_number != '*' else slide_number, tag_name

    def _find_shapes(self,
                     slide_number: Union[int, str],
                     tag_name: str) -> Iterable[BaseShape]:
        """Finds all shapes that match the label

        Args:
            label (str): label of the placeholder (without curly braces)
        """
        matched_shapes = []

        def _find_shapes_in_slide(slide):
            return filter(lambda shape: shape.text == f'{{{tag_name}}}', slide.shapes)

        if slide_number == '*':
            slides = self._presentation.slides
        else:
            # in label we are using 1 based indexing
            slide_index = slide_number - 1
            if slide_index < 0 or slide_index >= len(self._presentation.slides):
                raise IndexError(f"Can't find slide number {slide_number}.")
            slides = [self._presentation.slides[slide_index]]

        for slide in slides:
            matched_shapes.extend(_find_shapes_in_slide(slide))

        return matched_shapes

    def save(self, filename: _Pathlike) -> None:
        """Saves the updated pptx to the specified filepath.

        Args:
            filename (path-like): file name or path
        """
        # TODO: make sure that the user does not override the self._template_path
        pass
